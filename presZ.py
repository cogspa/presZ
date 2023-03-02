from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_HORIZONTAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# create a new PowerPoint presentation
prs = Presentation()

# add a new slide for each image
background_image_paths = [
    'path/to/loading_spinner_image.jpg',
    'path/to/hourglass_image.jpg',
    'path/to/person_with_loading_screen_image.jpg',
    'path/to/fun_fact_loading_screen_image.jpg',
    'path/to/branded_loading_screen_image.jpg'
]
for path in background_image_paths:
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # set the background image of the slide
    background_image = slide.shapes.add_picture(path, left=0, top=0, height=prs.slide_height, width=prs.slide_width)

    # create a white text box in the center of the slide
    left = Inches(1)
    top = Inches(1)
    width = prs.slide_width - Inches(2)
    height = prs.slide_height - Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
    textbox.line.color.rgb = RGBColor(255, 255, 255)
    textbox.line.width = Pt(0)
    textbox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    textbox.text_frame.horizontal_anchor = MSO_HORIZONTAL_ANCHOR.CENTER

    # add the bullet points to the text box
    bullet_points = [
        'Loading screens help manage user expectations by providing feedback during periods of waiting.',
        'They can help prevent user frustration and confusion by indicating that the app or website is still processing their request.',
        'A well-designed loading screen can serve as an opportunity to reinforce branding and create a positive impression on the user.',
        'Loading screens can also provide important information to the user, such as tips or interesting facts, to keep them engaged during the wait.',
        'Additionally, loading screens can be used strategically to distract the user from longer loading times by providing entertaining or informative content.'
    ]
    for point in bullet_points:
        paragraph = textbox.text_frame.add_paragraph()
        paragraph.text = point
        paragraph.font.size = Pt(60)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = MSO_ANCHOR.CENTER

# save the PowerPoint presentation
prs.save('presentation.pptx')
